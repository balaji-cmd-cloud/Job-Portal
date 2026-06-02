from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
PRIMARY_COLOR = RGBColor(79, 70, 229)      # Indigo
SECONDARY_COLOR = RGBColor(59, 130, 246)   # Blue
ACCENT_COLOR = RGBColor(236, 72, 153)      # Pink
DARK_TEXT = RGBColor(17, 24, 39)           # Dark Gray
LIGHT_BG = RGBColor(249, 250, 251)         # Light Gray
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = LIGHT_BG
    
    return slide

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    
    # Decorative line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(2), Inches(0))
    line.line.color.rgb = SECONDARY_COLOR
    line.line.width = Pt(3)
    
    # Content
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "APNA", "Production-Ready Job Portal with AI Capabilities")

# Slide 2: Project Overview
overview_items = [
    "🎯 Modern job portal platform inspired by Naukri",
    "🤖 AI-powered job recommendations & career chatbot",
    "💼 Built with Django, Bootstrap 5 & glassmorphism design",
    "🌙 Light/Dark theme support with responsive design",
    "📱 Full mobile & desktop compatibility",
    "⚡ REST API for seamless integration"
]
add_content_slide(prs, "Project Overview", overview_items)

# Slide 3: Key Features Part 1
features_1 = [
    "✉️  Email-Based Authentication",
    "   • Custom user model with email as primary key",
    "   • Case-insensitive email matching",
    "   • Strong password validation (8+ chars, uppercase, numbers, special)",
    "",
    "📊 Interactive Candidate Dashboard",
    "   • Profile progress tracking",
    "   • Saved & applied job statistics",
    "   • Dynamic job search filtering by domain, title & location"
]
add_content_slide(prs, "Key Features - Part 1", features_1)

# Slide 4: Key Features Part 2
features_2 = [
    "💾 Jobs Directory & Applications",
    "   • Full search listing with circular pagination",
    "   • AJAX-based save/apply without page reloads",
    "   • Similar job suggestions",
    "",
    "🤖 AI Chatbot",
    "   • Career guidance & resume tips",
    "   • Smart keyword-based counseling",
    "   • Conversation persistence in database",
    "   • Animated typing indicator"
]
add_content_slide(prs, "Key Features - Part 2", features_2)

# Slide 5: Key Features Part 3
features_3 = [
    "🧠 AI Recommendation Engine",
    "   • Skill coverage ratio calculation",
    "   • Experience level matching",
    "   • Location & domain preference matching",
    "   • Real-time trending skills aggregation",
    "",
    "🔌 REST API Layer",
    "   • Django REST Framework endpoints",
    "   • Jobs, Applications, Recommendations, Trending Skills"
]
add_content_slide(prs, "Key Features - Part 3", features_3)

# Slide 6: Tech Stack
tech_stack = [
    "Backend Framework: Django (Python)",
    "Frontend: Bootstrap 5, HTML, CSS, JavaScript",
    "Database: SQLite",
    "API: Django REST Framework (DRF)",
    "Design System: Glassmorphism with CSS Variables",
    "Theming: Light/Dark mode support",
    "UI Enhancements: AJAX, animations, responsive layout"
]
add_content_slide(prs, "Technology Stack", tech_stack)

# Slide 7: Project Structure
structure = [
    "job_portal/  → Main Django settings & URL routing",
    "accounts/    → Custom user authentication & user model",
    "jobs/        → Job listings, recommendations, API views",
    "chatbot/     → AI chatbot logic, conversation history",
    "static/      → CSS (glassmorphism), JavaScript (AJAX, theme)",
    "templates/   → HTML base structure",
    "media/       → User-uploaded files",
    "db.sqlite3   → Database"
]
add_content_slide(prs, "Project Structure", structure)

# Slide 8: API Endpoints
api_endpoints = [
    "GET  /jobs/api/                    → List all jobs (with filters)",
    "GET  /jobs/api/<id>/               → Get job details",
    "POST /jobs/api/<id>/save/          → Save/bookmark a job",
    "POST /jobs/api/<id>/apply/         → Apply for a job",
    "GET  /jobs/api/recommendations/    → Get personalized recommendations",
    "GET  /jobs/api/trending-skills/    → Get trending skills",
    "POST /chatbot/api/                 → Chat with AI assistant"
]
add_content_slide(prs, "REST API Endpoints", api_endpoints)

# Slide 9: Installation Steps
installation = [
    "1. Create Virtual Environment: python -m venv .venv",
    "",
    "2. Activate Virtual Environment (.venv/Scripts/activate)",
    "",
    "3. Install Dependencies: pip install -r requirements.txt",
    "",
    "4. Apply Migrations: python manage.py migrate",
    "",
    "5. Seed Database: python manage.py seed_jobs",
    "",
    "6. Create Superuser: python manage.py createsuperuser",
    "",
    "7. Run Server: python manage.py runserver"
]
add_content_slide(prs, "Installation & Setup", installation)

# Slide 10: Core Models
models = [
    "CustomUser",
    "   • Email-based authentication, profile data",
    "",
    "Job",
    "   • Job title, company, skills, salary, location, domain",
    "",
    "Application",
    "   • Tracks candidate applications & cover letters",
    "",
    "SavedJob / ChatHistory",
    "   • Bookmarks & conversation logs"
]
add_content_slide(prs, "Core Database Models", models)

# Slide 11: User Workflows
workflows = [
    "👤 Candidate User Flow:",
    "   Sign Up → Complete Profile → Search Jobs → Save/Apply → Track Applications",
    "",
    "💬 Chatbot Interaction:",
    "   Ask Career Question → AI Processes Keywords → Personalized Response",
    "",
    "🏢 Job Recommendations:",
    "   System Analyzes Profile → Calculates Match Score → Ranks Jobs",
    "",
    "📧 Email Authentication:",
    "   Register with Email → Validate Password → Access Dashboard"
]
add_content_slide(prs, "User Workflows", workflows)

# Slide 12: Design System
design = [
    "🎨 Glassmorphism Design",
    "   • Frosted glass effect with backdrop blur",
    "   • Transparent overlays with modern aesthetics",
    "",
    "🌈 Color Palette",
    "   • Primary: Indigo | Secondary: Blue | Accent: Pink",
    "",
    "📱 Responsive Breakpoints",
    "   • Mobile-first approach with Bootstrap 5",
    "   • Desktop, tablet, and mobile optimization",
    "",
    "⚫ Theme Support",
    "   • Light theme: Clean, bright interface",
    "   • Dark theme: Reduced eye strain, modern feel"
]
add_content_slide(prs, "Design System & UI", design)

# Slide 13: Key Algorithms
algorithms = [
    "🧮 Job Recommendation Algorithm",
    "   Skill Coverage Ratio + Experience Match + Location/Domain Bonus",
    "",
    "📊 Trending Skills Engine",
    "   Aggregates all job skill requirements in real-time",
    "",
    "🔑 Chatbot NLP Engine",
    "   Keyword extraction → Intent matching → Response template selection",
    "",
    "🔐 Custom Auth Backend",
    "   Case-insensitive email validation with Django ORM"
]
add_content_slide(prs, "Key Algorithms", algorithms)

# Slide 14: Testing & Quality
testing = [
    "✅ Unit Tests Available",
    "   • Authentication & registration validation",
    "   • Recommender algorithm scoring",
    "   • Chatbot history logging",
    "   • REST API endpoint coverage",
    "",
    "🧪 Test Command: python manage.py test",
    "",
    "📝 Admin Interface",
    "   • Django admin panel for data management"
]
add_content_slide(prs, "Testing & Quality Assurance", testing)

# Slide 15: Future Enhancements
enhancements = [
    "🚀 Planned Improvements",
    "",
    "• Advanced NLP using ML models (BERT, GPT)",
    "• Video interview integration",
    "• Resume parsing & skill extraction",
    "• Multi-language support",
    "• Payment gateway for premium features",
    "• Email notifications & job alerts",
    "• Social media profile integration"
]
add_content_slide(prs, "Future Enhancements", enhancements)

# Slide 16: Deployment
deployment = [
    "🌐 Production Deployment Options",
    "",
    "• Docker containerization",
    "• AWS/Azure cloud hosting",
    "• PostgreSQL for production database",
    "• Gunicorn/uWSGI WSGI server",
    "• Nginx reverse proxy",
    "• Redis caching layer",
    "• CI/CD pipeline with GitHub Actions"
]
add_content_slide(prs, "Deployment Strategy", deployment)

# Slide 17: Conclusion
conclusion = [
    "✨ APNA: A Complete Job Portal Solution",
    "",
    "🎯 Production-ready codebase with professional UI/UX",
    "🤖 AI-powered features for smart job matching",
    "📱 Fully responsive and accessible design",
    "🔌 RESTful API for modern integrations",
    "🧪 Comprehensive testing framework",
    "🚀 Scalable architecture ready for deployment"
]
add_content_slide(prs, "Conclusion", conclusion)

# Slide 18: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = SECONDARY_COLOR

thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
thanks_frame = thanks_box.text_frame
thanks_frame.word_wrap = True
thanks_p = thanks_frame.paragraphs[0]
thanks_p.text = "Thank You!"
thanks_p.font.size = Pt(60)
thanks_p.font.bold = True
thanks_p.font.color.rgb = WHITE
thanks_p.alignment = PP_ALIGN.CENTER

contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
contact_frame = contact_box.text_frame
contact_frame.word_wrap = True
contact_p = contact_frame.paragraphs[0]
contact_p.text = "For more information, visit the project repository or documentation"
contact_p.font.size = Pt(20)
contact_p.font.color.rgb = LIGHT_BG
contact_p.alignment = PP_ALIGN.CENTER

# Save presentation
prs.save('APNA_Project_Presentation.pptx')
print("✅ Presentation created: APNA_Project_Presentation.pptx")
